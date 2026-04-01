# Author : P.P. Chanchal
import sys
import asyncio
import os
import re
import json
import logging
import time
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any
import numpy as np
import pandas as pd
from PIL import Image
import torch
from sentence_transformers import SentenceTransformer
import emoji
import textract
from collections import Counter

# Robust imports for optional libraries
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] [AI_TXT] %(message)s'
)

class TextMetadata:
    """Data structure for text metadata"""
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.file_size = os.path.getsize(file_path)
        self.creation_time = os.path.getctime(file_path)
        self.modification_time = os.path.getmtime(file_path)
        self.word_count = 0
        self.char_count = 0
        self.language = "unknown"
        self.keywords = []
        self.sentiment_score = 0.0
        self.entities = []

    def to_dict(self):
        return {
            "file_path": self.file_path,
            "file_size": self.file_size,
            "created": self.creation_time,
            "modified": self.modification_time,
            "stats": {
                "words": self.word_count,
                "chars": self.char_count
            },
            "analysis": {
                "language": self.language,
                "keywords": self.keywords,
                "sentiment": self.sentiment_score,
                "entities": self.entities
            }
        }

class AdvancedTextProcessor:
    """
    Robust text processor capable of handling multiple formats,
    extracting metadata, and performing semantic analysis.
    """
    def __init__(self, model_name='all-MiniLM-L6-v2', device='cpu'):
        self.device = device
        logging.info(f"Loading SentenceTransformer: {model_name} on {device}")
        self.model = SentenceTransformer(model_name, device=device)
        self.emoji_cache = {}
        
        # Regex patterns
        self.email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
        self.url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
        self.phone_pattern = r'\+?[\d\s-]{10,}'

    async def extract_vector(self, file_path: Path) -> np.ndarray:
        """Main entry point for vector extraction"""
        text, metadata = await self.process_document(file_path)
        if not text:
            logging.warning(f"No text extracted from {file_path}")
            return np.zeros(self.model.get_sentence_embedding_dimension())
            
        # Enrich metadata
        metadata.word_count = len(text.split())
        metadata.char_count = len(text)
        metadata.keywords = self._extract_keywords(text)
        
        # Vectorize
        vector = self.model.encode(text, convert_to_numpy=True)
        return vector

    async def process_document(self, file_path: Path) -> Tuple[str, TextMetadata]:
        """Extracts text from various document formats"""
        metadata = TextMetadata(str(file_path))
        ext = file_path.suffix.lower()
        text = ""
        
        try:
            if ext == '.txt':
                text = self._read_txt(file_path)
            elif ext == '.pdf':
                text = self._read_pdf(file_path)
            elif ext == '.docx':
                text = self._read_docx(file_path)
            elif ext == '.pptx':
                text = self._read_pptx(file_path)
            elif ext in ['.csv', '.xlsx']:
                text = self._read_tabular(file_path)
            else:
                # Fallback to textract
                text = str(textract.process(str(file_path)).decode('utf-8'))
                
            # Clean text
            text = self._clean_text(text)
            return text, metadata
            
        except Exception as e:
            logging.error(f"Failed to process {file_path}: {e}")
            return "", metadata

    def _read_txt(self, path: Path) -> str:
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(path, 'r', encoding='latin-1') as f:
                return f.read()

    def _read_pdf(self, path: Path) -> str:
        if not fitz:
            logging.warning("PyMuPDF not installed, skipping PDF")
            return ""
        text = []
        with fitz.open(path) as doc:
            for page in doc:
                text.append(page.get_text())
        return "\n".join(text)

    def _read_docx(self, path: Path) -> str:
        if not docx:
            logging.warning("python-docx not installed, skipping DOCX")
            return ""
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs])

    def _read_pptx(self, path: Path) -> str:
        if not Presentation:
            logging.warning("python-pptx not installed, skipping PPTX")
            return ""
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _read_tabular(self, path: Path) -> str:
        try:
            if path.suffix == '.csv':
                df = pd.read_csv(path)
            else:
                df = pd.read_excel(path)
            return df.to_string()
        except Exception as e:
            logging.error(f"Tabular read error: {e}")
            return ""

    def _clean_text(self, text: str) -> str:
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove non-printable characters
        text = "".join(ch for ch in text if ch.isprintable())
        return text.strip()

    def _extract_keywords(self, text: str, top_k=5) -> List[str]:
        # Simple frequency-based keyword extraction
        # In a real system, use TF-IDF or RAKE
        words = re.findall(r'\w+', text.lower())
        stopwords = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'is', 'are', 'was', 'were'}
        filtered = [w for w in words if w not in stopwords and len(w) > 3]
        counter = Counter(filtered)
        return [word for word, _ in counter.most_common(top_k)]

    async def _process_text(self, text: str) -> np.ndarray:
        """Direct text vectorization helper"""
        return self.model.encode(text, convert_to_numpy=True)

    # --- Advanced Features (Placeholders for future expansion) ---
    def analyze_sentiment(self, text: str) -> float:
        # Placeholder: Return neutral
        return 0.0

    def extract_entities(self, text: str) -> List[Tuple[str, str]]:
        # Placeholder: Return empty
        return []

# Wrapper for compatibility
class UniversalFileProcessor(AdvancedTextProcessor):
    pass

if __name__ == "__main__":
    # Test
    async def test():
        proc = AdvancedTextProcessor()
        # Create dummy file
        with open("test_doc.txt", "w") as f:
            f.write("This is a test document for the autonomous AI system. It contains keywords like autonomy, robustness, and plasticity.")
        
        vec = await proc.extract_vector(Path("test_doc.txt"))
        print(f"Vector shape: {vec.shape}")
        
    asyncio.run(test())
